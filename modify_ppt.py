from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import os

src = "无线通信原理与移动网络 PPT.pptx"
dst = "无线通信原理与移动网络 PPT_修改版.pptx"

prs = Presentation(src)

def replace_text_in_slide(slide, old_text, new_text):
    """Replace text in all shapes of a slide."""
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join([run.text for run in paragraph.runs])
            if old_text in full_text:
                # Clear and rewrite preserving first run formatting
                if paragraph.runs:
                    first_run = paragraph.runs[0]
                    font_name = first_run.font.name
                    font_size = first_run.font.size
                    font_bold = first_run.font.bold
                    font_color = first_run.font.color.rgb if first_run.font.color.type else None
                else:
                    font_name, font_size, font_bold, font_color = None, None, None, None
                
                paragraph.clear()
                run = paragraph.add_run()
                run.text = full_text.replace(old_text, new_text)
                if font_name:
                    run.font.name = font_name
                if font_size:
                    run.font.size = font_size
                if font_bold is not None:
                    run.font.bold = font_bold
                if font_color:
                    run.font.color.rgb = font_color
                return True
    return False

def add_text_to_shape(shape, text):
    """Set text of a shape, trying to preserve formatting."""
    if hasattr(shape, "text_frame"):
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        return True
    return False

# Slide 10: replace placeholder
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text_frame.text
        if "请在此输入文字说明" in text:
            # This shape has the placeholder; replace entire content
            new_text = (
                "正交化可组合架构：网络结构（Standard/GNN/Hyper）与算法（IPPO/ExplabOff）完全解耦，\n"
                "通过统一 PolicyNetwork 接口实现任意组合，便于快速对比不同架构的泛化能力。"
            )
            add_text_to_shape(shape, new_text)
            break

# Slide 9: State -> Observation
slide9 = prs.slides[8]
replace_text_in_slide(slide9, "观测空间（State）", "观测空间（Observation）")

# Slide 8: fix latency description
slide8 = prs.slides[7]
replace_text_in_slide(
    slide8,
    "令同一时隙内本地执行与边缘执行并行，取两者延迟的最大值为最终延迟。若最终延迟大于1秒，则本次任务失败。",
    "任务延迟计算方法：若选择本地执行，延迟为本地计算时间；若选择边缘卸载，延迟为传输+排队+执行时间。若最终延迟大于1秒，则本次任务失败。"
)

# Slide 14: innovation overview
slide14 = prs.slides[13]
for shape in slide14.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text_frame.text
        if "关键词" in text or "正交架构" in text or "跨配置泛化" in text:
            new_text = (
                "本项目的核心技术创新：\n"
                "1. 正交化可组合架构：网络结构与训练算法解耦，Standard/GNN/Hyper × IPPO/ExplabOff 任意组合；\n"
                "2. 跨配置网络架构：以 GNN 和 HyperNetwork 替代论文基线 MLP，单模型服务多 MEC 配置；\n"
                "3. 训练稳定性工程：LR decay、Trajectory buffer 管理等优化，保障长时训练稳定收敛。"
            )
            add_text_to_shape(shape, new_text)
            break

# Slide 15: network architecture optimization (main innovation)
slide15 = prs.slides[14]
# Find largest text-frame shape (likely body placeholder)
body15 = None
for shape in slide15.shapes:
    if hasattr(shape, "text_frame"):
        if body15 is None or (shape.width * shape.height) > (body15.width * body15.height):
            body15 = shape
if body15 is not None:
    new_text = (
        "从论文 MLP 到可扩展网络架构\n\n"
        "论文基线：\n"
        "• 每个 (M, E) 配置独立训练一个 MLP；\n"
        "• 输入维度固定为 obs_dim = 1 + E，无法跨配置复用；\n"
        "• 新增配置必须重新训练、重新调参。\n\n"
        "我们的优化：\n"
        "• GNN：图注意力网络，MD/ES 作为节点共享参数，天然支持任意节点数；\n"
        "• HyperNetwork：配置编码器动态生成策略网络权重，(M, E) 作为输入；\n"
        "• 统一 PolicyNetwork 接口：三种网络对训练算法完全透明，即插即用。"
    )
    add_text_to_shape(body15, new_text)

# Slide 16: training stability engineering
slide16 = prs.slides[15]
body16 = None
for shape in slide16.shapes:
    if hasattr(shape, "text_frame"):
        if body16 is None or (shape.width * shape.height) > (body16.width * body16.height):
            body16 = shape
if body16 is not None:
    new_text = (
        "训练稳定性优化\n\n"
        "• 学习率衰减（LR Decay）：每 5K episodes 降低 0.5 倍，避免后期震荡；\n"
        "• Trajectory Buffer 管理：每次 PPO update 后清空 buffer，防止数据累积；\n"
        "• GNN 过平滑修复：1-layer GNN + agent embedding，解决所有 MD 决策同质化；\n"
        "• Action Masking（可选）：根据任务大小与 ES 容量预屏蔽无效动作。"
    )
    add_text_to_shape(body16, new_text)

# Slide 19: update key findings
slide19 = prs.slides[18]
# We need to replace the entire findings list. Find the shape containing the findings.
for shape in slide19.shapes:
    if hasattr(shape, "text_frame"):
        text = shape.text_frame.text
        if "IPPO 在复杂环境胜出" in text or "关键发现" in text:
            new_text = (
                "关键发现：\n"
                "1. ExplabOff+GNN 泛化最优：在3ES-7MD未见数据上，ExplabOff+GNN（0.412）比 IPPO+GNN（0.461）好 10.6%；\n"
                "2. 训练最优 ≠ 泛化最优：IPPO+Hyper 训练峰值最低（0.400），但评估时退化严重（0.526）；\n"
                "3. GNN 过平滑修复：1-layer GNN + agent embedding 将 2ES-3MD cost 从 0.998 降至 0.426（-57%）；\n"
                "4. GNN 是唯一真正跨配置架构：Standard MLP 受固定 obs_dim 限制，2ES↔3ES 不兼容；\n"
                "5. Trajectory buffer 清空：忘记清空导致数据累积 100K+，训练失效。"
            )
            add_text_to_shape(shape, new_text)
            break

# Save modified presentation
prs.save(dst)
print(f"Modified PPT saved to: {dst}")
print(f"File size: {os.path.getsize(dst) / 1024:.1f} KB")
